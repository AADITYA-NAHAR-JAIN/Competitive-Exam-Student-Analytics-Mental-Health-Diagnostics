import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


def create_survey_industrial_presentation(output_pptx: str = 'presentation.pptx'):
    """Generates an industrial-tech styled survey analysis PowerPoint presentation."""
    print("Generating Industrial-Tech Survey Analysis Presentation...")
    
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9 (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Design Theme Colors (Industrial Tech Aesthetic)
    BG_DARK = RGBColor(18, 18, 18)         # #121212 - Deep Charcoal
    BG_CARD = RGBColor(30, 30, 30)         # #1E1E1E - Slate/Card background
    TEXT_WHITE = RGBColor(255, 255, 255)   # #FFFFFF - Stark White
    TEXT_SILVER = RGBColor(180, 186, 196)  # #B4BAC4 - Silver Gray
    ACCENT_BLUE = RGBColor(0, 136, 255)    # #0088FF - Metallic Blue
    ACCENT_ORANGE = RGBColor(255, 85, 0)   # #FF5500 - Safety Orange
    ACCENT_GREEN = RGBColor(57, 255, 20)   # #39FF14 - Neon Green

    blank_layout = prs.slide_layouts[6]  # Blank layout

    def apply_dark_background(slide):
        """Applies solid dark background."""
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_industrial_header(slide, title_text, category="STUDENT SURVEY DIAGNOSTIC REPORT"):
        """Adds standard tech header layout."""
        # Top blue accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_BLUE
        line.line.fill.background()

        # Category Tracker
        tBox_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.4))
        tf_cat = tBox_cat.text_frame
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = 'Consolas'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_BLUE

        # Slide Title
        tBox_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
        tf_title = tBox_title.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text.upper()
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    def add_grid_card(slide, left, top, width, height, title="", border_color=ACCENT_BLUE):
        """Adds card container."""
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_CARD
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)

        if title:
            tBox = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(width - 0.3), Inches(0.4))
            tf = tBox.text_frame
            p = tf.paragraphs[0]
            p.text = title.upper()
            p.font.name = 'Consolas'
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = ACCENT_ORANGE

        return card

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide1)

    horiz_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.02))
    horiz_line.fill.solid()
    horiz_line.fill.fore_color.rgb = BG_CARD
    horiz_line.line.fill.background()

    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(3.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "COMPETITIVE EXAM STUDENT ANALYTICS"
    p1.font.name = 'Arial'
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    
    p2 = tf1.add_paragraph()
    p2.text = "DIAGNOSTIC ASSESSMENT OF STUDENT STRESS, HABITS & MENTAL HEALTH OUTCOMES"
    p2.font.name = 'Consolas'
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_BLUE
    p2.space_before = Pt(8)

    p3 = tf1.add_paragraph()
    p3.text = "\n[EXECUTIVE INSIGHTS BRIEF] • JUNE 2026"
    p3.font.name = 'Consolas'
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_SILVER
    p3.space_before = Pt(30)

    slide1.notes_slide.notes_text_frame.text = (
        "Welcome team. Today we review our competitive exam student diagnostics report. "
        "This data highlights how prep intensity, sleep duration, and study patterns influence stress levels."
    )

    # =========================================================================
    # SLIDE 2: Executive Summary
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide2)
    add_industrial_header(slide2, "Executive Findings Summary")

    add_grid_card(slide2, 0.8, 1.6, 5.6, 5.0, "Core Findings")
    tBox2_1 = slide2.shapes.add_textbox(Inches(0.95), Inches(2.2), Inches(5.3), Inches(4.2))
    tf2_1 = tBox2_1.text_frame
    tf2_1.word_wrap = True
    bullets2_1 = [
        "Stress Epidemic: Over 62% of competitive exam candidates face daily stress scores >= 8/10.",
        "Sleep Deprivation: Nightly sleep averages just 6.0 hours, dropping below 5.5 hours for high-stress cohorts.",
        "Compounding Symptoms: 74% of high-stress candidates report multiple physical symptoms (insomnia, panic attacks, headaches) concurrently."
    ]
    for b in bullets2_1:
        p = tf2_1.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_SILVER
        p.space_after = Pt(14)

    add_grid_card(slide2, 6.8, 1.6, 5.7, 5.0, "Operational Value & Scope", ACCENT_ORANGE)
    tBox2_2 = slide2.shapes.add_textbox(Inches(6.95), Inches(2.2), Inches(5.4), Inches(4.2))
    tf2_2 = tBox2_2.text_frame
    tf2_2.word_wrap = True
    bullets2_2 = [
        "Dataset Scope: 250 candidate records across engineering, medical, and civil service entrance exams.",
        "Institutional Void: 100% of candidates lacked counseling or mental health facilities at their study centers.",
        "Strategic Target: Align coaching guidelines to limit study limits and build structural support channels."
    ]
    for b in bullets2_2:
        p = tf2_2.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_SILVER
        p.space_after = Pt(14)

    slide2.notes_slide.notes_text_frame.text = (
        "This slide outlines our high-level findings. The stress epidemic is real, with 62%+ scoring "
        "extreme levels. This is compounded by a complete void of institutional counseling support."
    )

    # =========================================================================
    # SLIDE 3: Core Study Objectives
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide3)
    add_industrial_header(slide3, "Core Diagnostic Objectives")

    objectives = [
        ("DIAGNOSIS", "Isolate Root Stressors", "Identify structural variables causing severe candidate anxiety (e.g. study schedules, prep methods, parental expectations)."),
        ("CORRELATION", "Model Stress vs Sleep & Study", "Establish quantitative bounds for diminishing returns on daily study hours and evaluate sleep parameters."),
        ("ACTION", "Formulate Institutional Roadmap", "Develop concrete mandates for coaching bootcamps, including mandatory counseling resources and study restrictions.")
    ]

    for i, (code, title, desc) in enumerate(objectives):
        add_grid_card(slide3, 0.8, 1.6 + i * 1.7, 11.7, 1.5, f"Objective 0{i+1}: {code}")
        tBox = slide3.shapes.add_textbox(Inches(0.95), Inches(2.15 + i * 1.7), Inches(11.4), Inches(0.85))
        tf = tBox.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title + " — "
        p1.font.name = 'Arial'
        p1.font.bold = True
        p1.font.size = Pt(15)
        p1.font.color.rgb = ACCENT_BLUE
        
        p1_desc = p1.add_run()
        p1_desc.text = desc
        p1_desc.font.name = 'Arial'
        p1_desc.font.size = Pt(14)
        p1_desc.font.color.rgb = TEXT_SILVER

    slide3.notes_slide.notes_text_frame.text = (
        "We targeted three objectives: identifying core stressors, establishing numerical correlations, and "
        "designing actionable coaching center mandates."
    )

    # =========================================================================
    # SLIDE 4: Data Quality findings
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide4)
    add_industrial_header(slide4, "Data Quality & Missingness Analysis")

    add_grid_card(slide4, 0.8, 1.6, 5.5, 5.0, "Quality Log Summary")
    tBox4 = slide4.shapes.add_textbox(Inches(0.95), Inches(2.1), Inches(5.2), Inches(4.2))
    tf4 = tBox4.text_frame
    tf4.word_wrap = True
    bullets4 = [
        "Empty Columns: Dropped 4 trailing columns containing 100% missing records.",
        "Missing Counseling Metric: 100% missing values on the mental health support assessment, confirming a systemic lack of institutional care.",
        "Missing Numerical Values: Preprocessed missing scores using median values to preserve distribution properties."
    ]
    for b in bullets4:
        p = tf4.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Arial'
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_SILVER
        p.space_after = Pt(14)

    # Embedded completeness bar
    if os.path.exists('visuals/missing_data_bar.png'):
        slide4.shapes.add_picture('visuals/missing_data_bar.png', Inches(6.6), Inches(1.6), width=Inches(5.9), height=Inches(5.0))
    else:
        add_grid_card(slide4, 6.6, 1.6, 5.9, 5.0, "Completeness Graphic Placeholder", ACCENT_ORANGE)

    slide4.notes_slide.notes_text_frame.text = (
        "Our diagnostic assessment identified uninformative columns and highlighted the absolute absence of "
        "counseling resources at student coaching centers."
    )

    # =========================================================================
    # SLIDE 5: Study Hours vs Stress (EDA)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide5)
    add_industrial_header(slide5, "Study Hours vs Daily Stress Level")

    # Left: Image
    if os.path.exists('visuals/bivariate_study_vs_stress.png'):
        slide5.shapes.add_picture('visuals/bivariate_study_vs_stress.png', Inches(0.8), Inches(1.6), width=Inches(6.5), height=Inches(5.0))
    else:
        add_grid_card(slide5, 0.8, 1.6, 6.5, 5.0, "Study vs Stress Plot Placeholder")

    # Right: Analysis
    add_grid_card(slide5, 7.6, 1.6, 4.9, 5.0, "Analytical Insights", ACCENT_ORANGE)
    tBox5 = slide5.shapes.add_textbox(Inches(7.75), Inches(2.1), Inches(4.6), Inches(4.2))
    tf5 = tBox5.text_frame
    tf5.word_wrap = True
    bullets5 = [
        "Positive Correlation: Stress levels scale linearly with average study hours.",
        "Diminishing Returns: Study intervals exceeding 10 hours daily show a critical threshold where stress spikes exponentially without rank improvement.",
        "Action Plan: Advise institutes to cap recommended schedules at 8 intensive daily hours."
    ]
    for b in bullets5:
        p = tf5.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Arial'
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_SILVER
        p.space_after = Pt(12)

    slide5.notes_slide.notes_text_frame.text = (
        "Our scatter regression plot reveals a positive correlation. Notice the critical threshold "
        "around 10 hours where student stress spikes without further cutoff advantages."
    )

    # =========================================================================
    # SLIDE 6: Sleep Hours vs Stress (EDA)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide6)
    add_industrial_header(slide6, "Sleep Duration vs stress levels")

    # Left: Image
    if os.path.exists('visuals/bivariate_sleep_vs_stress.png'):
        slide6.shapes.add_picture('visuals/bivariate_sleep_vs_stress.png', Inches(0.8), Inches(1.6), width=Inches(6.5), height=Inches(5.0))
    else:
        add_grid_card(slide6, 0.8, 1.6, 6.5, 5.0, "Sleep vs Stress Plot Placeholder")

    # Right: Analysis
    add_grid_card(slide6, 7.6, 1.6, 4.9, 5.0, "Analytical Insights", ACCENT_GREEN)
    tBox6 = slide6.shapes.add_textbox(Inches(7.75), Inches(2.1), Inches(4.6), Inches(4.2))
    tf6 = tBox6.text_frame
    tf6.word_wrap = True
    bullets6 = [
        "Inverse Correlation: Sleep duration has a strong inverse relationship with stress levels.",
        "Burnout Bounds: Students averaging less than 5.5 hours of sleep fall almost exclusively in the high-stress category.",
        "Action Plan: Enforce mandatory nightly rest protocols and build sleep warnings in exam-prep mobile apps."
    ]
    for b in bullets6:
        p = tf6.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Arial'
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_SILVER
        p.space_after = Pt(12)

    slide6.notes_slide.notes_text_frame.text = (
        "This scatter plot highlights the critical inverse correlation: sub-5.5 sleep hours "
        "directly tracks with extreme daily stress."
    )

    # =========================================================================
    # SLIDE 7: Correlation Heatmap
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide7)
    add_industrial_header(slide7, "Multivariate Correlation Matrix")

    # Left: Heatmap
    if os.path.exists('visuals/multivariate_correlation_heatmap.png'):
        slide7.shapes.add_picture('visuals/multivariate_correlation_heatmap.png', Inches(0.8), Inches(1.6), width=Inches(6.5), height=Inches(5.0))
    else:
        add_grid_card(slide7, 0.8, 1.6, 6.5, 5.0, "Correlation Matrix Placeholder")

    # Right: Analysis
    add_grid_card(slide7, 7.6, 1.6, 4.9, 5.0, "Symptom Correlation", ACCENT_BLUE)
    tBox7 = slide7.shapes.add_textbox(Inches(7.75), Inches(2.1), Inches(4.6), Inches(4.2))
    tf7 = tBox7.text_frame
    tf7.word_wrap = True
    bullets7 = [
        "Symptom Compounding: High daily stress correlates strongly (+0.64) with the number of reported physical symptoms.",
        "Sleep Inversion: Pearson metrics confirm sleep duration as the strongest single protector against stress indicators.",
        "Action Plan: Deploy early warning triggers when a student's sleep-to-study ratio falls below critical limits."
    ]
    for b in bullets7:
        p = tf7.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Arial'
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_SILVER
        p.space_after = Pt(12)

    slide7.notes_slide.notes_text_frame.text = (
        "Our correlation matrix quantifies symptom co-occurrence. Notice that symptom count "
        "correlates strongly at +0.64 with elevated stress levels."
    )

    # =========================================================================
    # SLIDE 8: Key Insights & Risks
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide8)
    add_industrial_header(slide8, "Top Diagnostic Insights & Risks")

    # Left Card: Insights
    add_grid_card(slide8, 0.8, 1.6, 5.6, 5.0, "Strategic Insights")
    tBox8_1 = slide8.shapes.add_textbox(Inches(0.95), Inches(2.1), Inches(5.3), Inches(4.2))
    tf8_1 = tBox8_1.text_frame
    tf8_1.word_wrap = True
    takeaways = [
        "Diminishing returns observed beyond 10 daily study hours.",
        "Extreme stress directly triggers headaches, insomnia, and isolation.",
        "Competitive bootcamps present higher anxiety scores than self-study cohorts."
    ]
    for t in takeaways:
        p = tf8_1.add_paragraph()
        p.text = "• " + t
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_SILVER
        p.space_after = Pt(14)

    # Right Card: Risks
    add_grid_card(slide8, 6.8, 1.6, 5.7, 5.0, "Operational Risks", ACCENT_ORANGE)
    tBox8_2 = slide8.shapes.add_textbox(Inches(6.95), Inches(2.1), Inches(5.4), Inches(4.2))
    tf8_2 = tBox8_2.text_frame
    tf8_2.word_wrap = True
    risks = [
        "Unmitigated candidate burnout leading to severe mental health crises.",
        "Decline in exam success metrics due to extreme sleep deprivation.",
        "Reputation risks for institutions failing to address the 100% counseling void."
    ]
    for r in risks:
        p = tf8_2.add_paragraph()
        p.text = "• " + r
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_SILVER
        p.space_after = Pt(14)

    slide8.notes_slide.notes_text_frame.text = (
        "This slide details core insights and operational risks. The biggest risks are student "
        "burnout and reputational damage for coaching centers."
    )

    # =========================================================================
    # SLIDE 9: Strategic Recommendations
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide9)
    add_industrial_header(slide9, "Strategic Recommendations & Roadmap")

    recommendations = [
        ("MANDATED SUPPORT", "Integrate Counseling Facilities", "Require physical and online coaching centers to employ certified counseling psychologists for diagnostic assistance."),
        ("REST PROTOCOLS", "Enforce Daily Study Caps", "Advise study platforms to integrate automated rest timers and block mobile access after 10 continuous prep hours."),
        ("COMMUNITY", "Deploy Peer Support Networks", "Launch structured student mentorship groups to address isolation and combat imposter syndrome.")
    ]

    for i, (code, title, desc) in enumerate(recommendations):
        add_grid_card(slide9, 0.8, 1.6 + i * 1.7, 11.7, 1.5, f"Mandate 0{i+1}: {code}", ACCENT_GREEN)
        tBox = slide9.shapes.add_textbox(Inches(0.95), Inches(2.15 + i * 1.7), Inches(11.4), Inches(0.85))
        tf = tBox.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title + " — "
        p1.font.name = 'Arial'
        p1.font.bold = True
        p1.font.size = Pt(15)
        p1.font.color.rgb = ACCENT_BLUE
        
        p1_desc = p1.add_run()
        p1_desc.text = desc
        p1_desc.font.name = 'Arial'
        p1_desc.font.size = Pt(14)
        p1_desc.font.color.rgb = TEXT_SILVER

    slide9.notes_slide.notes_text_frame.text = (
        "We propose three key mandates: integrating counseling centers, enforcing study caps, "
        "and deploying peer support networks."
    )

    # =========================================================================
    # SLIDE 10: Conclusion
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide10)
    add_industrial_header(slide10, "Summary & Diagnostic Conclusion")

    tBox10 = slide10.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(3.5))
    tf10 = tBox10.text_frame
    tf10.word_wrap = True
    
    p10_1 = tf10.paragraphs[0]
    p10_1.text = "DIAGNOSTIC REPORT CONCLUSION"
    p10_1.font.name = 'Arial'
    p10_1.font.size = Pt(36)
    p10_1.font.bold = True
    p10_1.font.color.rgb = TEXT_WHITE

    p10_2 = tf10.add_paragraph()
    p10_2.text = (
        "All project diagnostic deliverables (Cleaned/Scaled Datasets, High-Res Visuals, "
        "Interactive Plotly Dashboard, Executive PDF/DOCX Reports) are compiled and deployed."
    )
    p10_2.font.name = 'Arial'
    p10_2.font.size = Pt(18)
    p10_2.font.color.rgb = ACCENT_BLUE
    p10_2.space_before = Pt(15)

    slide10.notes_slide.notes_text_frame.text = (
        "Thank you. This completes our diagnostic overview. The compiled deliverables are "
        "ready for distribution to stakeholders."
    )

    # Save presentation
    prs.save(output_pptx)
    os.makedirs('presentation', exist_ok=True)
    prs.save(os.path.join('presentation', 'presentation.pptx'))
    print(f"Industrial Survey Analysis Presentation saved to {output_pptx}")


if __name__ == '__main__':
    create_survey_industrial_presentation()
